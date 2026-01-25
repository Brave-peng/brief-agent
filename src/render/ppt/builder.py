"""
直接渲染 PPT 构建器

通过 LLM 生成设计蓝图，直接使用 python-pptx 渲染
"""
import json
import logging
import os
import re
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from src.agents import get_llm
from .base import PPTBuilder, BuilderRegistry

logger = logging.getLogger(__name__)


@BuilderRegistry.register("direct")
class DirectPPBuilder(PPTBuilder):
    """直接渲染 PPT 构建器"""

    DEFAULT_STYLE = "干净简洁的插画风格，黑色线条艺术，纯白背景 --ar 16:9"

    PROMPT_TEMPLATE = """
你是一位世界顶级的视觉叙事艺术家。请将旁白文本转换成 PPT 幻灯片设计蓝图。

**设计原则:**
1. 视觉优先，大胆使用全屏图片
2. 信息降噪，只展示关键词或金句
3. 统一视觉风格: {style}

**输出 JSON:**
```json
{{
  "design_concept": "设计思路描述",
  "slide_title": {{"text": "标题", "position": {{"left": 0.5, "top": 0.5, "width": 9.0, "height": 1.0, "font_size": 32}}}},
  "content_elements": [{{"text": "核心要点", "position": {{"left": 1.0, "top": 2.0, "width": 4.0, "height": 1.5, "font_size": 18}}}}],
  "image_elements": [{{"prompt": "英文图片提示词", "position": {{"left": 5.5, "top": 1.5, "width": 4.0, "height": 5.0}}}}]
}}
```

幻灯片尺寸: 10 x 7.5 英寸
请为以下旁白设计幻灯片:
{script}
"""

    def __init__(
        self,
        data: dict[str, Any],
        provider: str = "minimax",
        style: str | None = None,
    ) -> None:
        self.data = data
        self.llm = get_llm()
        self.provider = provider
        self.style = style or self.DEFAULT_STYLE
        self.presentation = Presentation()
        self._blank_layout = self._get_blank_layout()

    def _get_blank_layout(self) -> Any:
        for layout in self.presentation.slide_layouts:
            if layout.name in ["Blank", "空白"]:
                return layout
        return self.presentation.slide_layouts[-1]

    def _sanitize_filename(self, text: str) -> str:
        if not text:
            return "untitled"
        return re.sub(r'[\\/*?:"<>|]', "", text)[:20]

    def _get_blueprint(self, script: str) -> dict[str, Any]:
        prompt = self.PROMPT_TEMPLATE.format(style=self.style, script=script)
        resp = self.llm.generate(prompt, provider=self.provider)
        resp = resp.strip().replace("```json", "").replace("```", "").strip()
        try:
            return json.loads(resp)
        except json.JSONDecodeError:
            logger.error("JSON 解析失败，使用空蓝图")
            return {"slide_title": {"text": "解析失败"}, "content_elements": [], "image_elements": []}

    def _render_slide(self, blueprint: dict[str, Any], slide_id: str) -> None:
        slide = self.presentation.slides.add_slide(self._blank_layout)

        # 图片
        for img in blueprint.get("image_elements", []):
            path = img.get("image_path_to_render")
            pos = img.get("position", {})
            if path and os.path.exists(path) and pos:
                slide.shapes.add_picture(
                    path,
                    Inches(pos.get("left", 0)),
                    Inches(pos.get("top", 0)),
                    width=Inches(pos.get("width", 4)),
                    height=Inches(pos.get("height", 3)),
                )

        # 标题
        title = blueprint.get("slide_title", {})
        if title.get("text"):
            pos = title.get("position", {})
            box = slide.shapes.add_textbox(
                Inches(pos.get("left", 0)),
                Inches(pos.get("top", 0)),
                Inches(pos.get("width", 10)),
                Inches(pos.get("height", 1)),
            )
            p = box.text_frame.paragraphs[0]
            p.text = title["text"]
            p.font.size = Pt(pos.get("font_size", 32))
            p.font.bold = True

        # 内容
        for elem in blueprint.get("content_elements", []):
            pos = elem.get("position", {})
            box = slide.shapes.add_textbox(
                Inches(pos.get("left", 0)),
                Inches(pos.get("top", 0)),
                Inches(pos.get("width", 5)),
                Inches(pos.get("height", 2)),
            )
            box.text_frame.text = elem.get("text", "")
            box.text_frame.paragraphs[0].font.size = Pt(pos.get("font_size", 18))

    def build(self, output_path: str) -> None:
        """构建 PPT"""
        slides = self.data.get("slides", [])

        # 标题页
        title = self.data.get("title", "日报")
        self._render_title_slide(title)

        # 内容页
        for i, slide in enumerate(slides):
            slide_title = slide.get("title", f"话题 {i+1}")
            key_points = slide.get("key_points", "")
            bullet_points = slide.get("bullet_points", [])
            speaker_notes = slide.get("speaker_notes", "")

            self._render_content_slide(slide_title, key_points, bullet_points)

            # 添加 speaker_notes 到 PPT 备注
            if speaker_notes:
                self._add_simple_notes(speaker_notes)

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(output_path)
        logger.info("PPT 已保存到: %s", output_path)

    def _render_title_slide(self, title: str) -> None:
        """渲染标题页"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        slide.shapes.title.text = title
        # 副标题放摘要
        if hasattr(slide.placeholders[1], "text"):
            slide.placeholders[1].text = "AI 驱动的技术日报"

    def _render_content_slide(self, title: str, key_points: str, bullet_points: list) -> None:
        """渲染内容页（固定布局）"""
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        slide.shapes.title.text = title

        # 关键点作为副标题
        if hasattr(slide.placeholders[1], "text"):
            slide.placeholders[1].text = key_points

        # 要点列表
        content = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if content and hasattr(content, "text_frame"):
            tf = content.text_frame
            tf.clear()
            for point in bullet_points:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(18)

    def _add_simple_notes(self, notes: str) -> None:
        """添加演讲者备注（简化版）"""
        slide = self.presentation.slides[-1]
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    def _add_speaker_notes(self, slide_id: str, notes: str) -> None:
        """添加演讲者备注"""
        # 查找对应的 slide
        slide_index = int(slide_id.replace("slide_", "")) if "slide_" in slide_id else 0
        if slide_index < len(self.presentation.slides):
            slide = self.presentation.slides[slide_index]
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
